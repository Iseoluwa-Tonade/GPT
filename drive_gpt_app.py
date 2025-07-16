import os
import streamlit as st
import json
import openai
from google_auth_oauthlib.flow import Flow
from googleapiclient.discovery import build
import io
import base64
import fitz  # PyMuPDF
import docx
import pptx
from PIL import Image
from google.oauth2.credentials import Credentials
from google.auth.transport.requests import Request
from googleapiclient.errors import HttpError
from googleapiclient.http import MediaIoBaseDownload

# --- Configuration ---
st.set_page_config(page_title="File Analyzer", page_icon="🧩")
st.title("🧩 Google Drive File Analyzer")
st.write("Select one or more files from Google Drive for analysis.")

# --- Constants ---
SCOPES = ["https://www.googleapis.com/auth/drive.readonly"]
# NOTE: The redirect URI for deployment must match the one in your Google Cloud console.
REDIRECT_URI = "https://zw2bm6uwryon2f5pnfsauk.streamlit.app/"

# --- Helper Functions ---
def get_file_content(drive_service, file_info):
    """Downloads and extracts content from a Google Drive file."""
    file_id = file_info['id']
    mime_type = file_info['mimeType']
    try:
        request = drive_service.files().get_media(fileId=file_id)
        file_bytes = io.BytesIO()
        downloader = MediaIoBaseDownload(file_bytes, request)
        done = False
        while not done:
            status, done = downloader.next_chunk()
        file_bytes.seek(0)
        
        if 'pdf' in mime_type:
            with fitz.open(stream=file_bytes, filetype="pdf") as doc:
                return "text", "".join(page.get_text() for page in doc)
        elif 'vnd.openxmlformats-officedocument.wordprocessingml.document' in mime_type:
            doc = docx.Document(file_bytes)
            return "text", "\n".join([para.text for para in doc.paragraphs])
        elif 'vnd.openxmlformats-officedocument.presentationml.presentation' in mime_type:
            prs = pptx.Presentation(file_bytes)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_runs.append(run.text)
            return "text", "\n".join(text_runs)
        elif 'image' in mime_type:
            # Verify image is valid before returning bytes
            Image.open(file_bytes)
            return "image", file_bytes.getvalue()
        elif 'text' in mime_type:
            return "text", file_bytes.getvalue().decode("utf-8", errors='ignore')
        else:
            return "unsupported", f"File type ('{mime_type}')"
    
    except HttpError:
        # Handle Google Workspace files (Docs, Sheets, Slides) which need to be exported
        if 'google-apps' in mime_type:
            if mime_type == 'application/vnd.google-apps.shortcut':
                return "unsupported", "Google Drive Shortcut"
            
            export_mime_type = None
            if mime_type == 'application/vnd.google-apps.document':
                export_mime_type = 'text/plain'
            elif mime_type == 'application/vnd.google-apps.spreadsheet':
                export_mime_type = 'text/csv'
            elif mime_type == 'application/vnd.google-apps.presentation':
                export_mime_type = 'text/plain'

            if export_mime_type:
                try:
                    request = drive_service.files().export_media(fileId=file_id, mimeType=export_mime_type)
                    file_bytes = io.BytesIO()
                    downloader = MediaIoBaseDownload(file_bytes, request)
                    done = False
                    while not done:
                        status, done = downloader.next_chunk()
                    return "text", file_bytes.getvalue().decode("utf-8")
                except HttpError as e:
                    return "unsupported", f"Export Error: {e}"
            else:
                return "unsupported", f"Google Workspace type ({mime_type})"
        else:
            return "unsupported", "Google Drive API Error"

def get_ai_response(api_key, model, messages, max_tokens=4000):
    """Generic function to call the OpenAI ChatCompletion API."""
    try:
        client = openai.OpenAI(api_key=api_key)
        response = client.chat.completions.create(
            model=model,
            messages=messages,
            max_tokens=max_tokens
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"ERROR: {str(e)}"

# --- Main App Logic ---
if "credentials" not in st.session_state:
    st.session_state.credentials = None

if st.session_state.credentials is None:
    st.subheader("Step 1: Authenticate with Google")
    try:
        # Load credentials from Streamlit secrets
        client_config = st.secrets["google_credentials"]
        flow = Flow.from_client_config(client_config, scopes=SCOPES, redirect_uri=REDIRECT_URI)
        
        auth_url, _ = flow.authorization_url(prompt="consent")
        st.link_button("Login with Google", auth_url, help="Click to authorize access to your Google Drive files. You will be redirected to a Google consent screen.")
        
        st.write("After authorizing, you will be redirected back to the app. If not, copy the full redirected URL from your browser and paste it below.")
        redirected_url = st.text_input("Paste the full redirected URL here:")
        
        if redirected_url:
            # Extract the authorization code from the URL
            code = redirected_url.split('code=')[1].split('&scope')[0]
            flow.fetch_token(code=code)
            creds = flow.credentials
            st.session_state.credentials = creds.to_json()
            st.rerun()

    except Exception as e:
        st.error("Could not load Google credentials from secrets. Ensure they are correctly configured in your Streamlit Cloud settings.")
        st.error(f"Specific error: {e}")
else:
    creds = Credentials.from_authorized_user_info(json.loads(st.session_state.credentials))
    # Refresh credentials if they have expired
    if creds.expired and creds.refresh_token:
        creds.refresh(Request())
        st.session_state.credentials = creds.to_json()

    drive_service = build("drive", "v3", credentials=creds)
    st.success("✅ Connected to Google Drive!")
    if st.button("Logout"):
        st.session_state.credentials = None
        st.rerun()
    st.divider()

    st.subheader("Step 2: Select Files and Ask a Question")
    try:
        # List files from Google Drive (excluding folders)
        results = drive_service.files().list(
            pageSize=200, 
            fields="files(id, name, mimeType)",
            q="mimeType != 'application/vnd.google-apps.folder'"
        ).execute()
        files = results.get("files", [])
        
        if not files:
            st.write("No files found in your Google Drive.")
        else:
            file_options = {f"{file['name']} ({file['mimeType']})": file for file in files}
            selected_files_display = st.multiselect(
                "Choose files to analyze:", 
                options=list(file_options.keys())
            )
            
            user_prompt = st.text_area("What would you like to know about these files?", height=100)

            if st.button("✨ Analyze Files", disabled=(not selected_files_display or not user_prompt)):
                api_key = st.secrets["openai"]["api_key"]
                summaries = []
                image_parts = []
                
                with st.spinner("Processing and summarizing files... This may take a moment."):
                    # --- MAP STEP: Summarize each document individually ---
                    for file_display in selected_files_display:
                        file_info = file_options[file_display]
                        st.info(f"Processing: {file_info['name']}")
                        content_type, content = get_file_content(drive_service, file_info)

                        if content_type == 'text':
                            # Create a simple summarization prompt for each file
                            summary_prompt = f"Summarize the key points of the following document named '{file_info['name']}':\n\n{content}"
                            summary_messages = [{"role": "user", "content": summary_prompt}]
                            
                            # Use a cheaper/faster model for summarization and request a smaller response.
                            summary = get_ai_response(api_key, model="gpt-4o-mini", messages=summary_messages, max_tokens=500)
                            
                            if summary.startswith("ERROR:"):
                                st.error(f"Could not summarize {file_info['name']}: {summary}")
                                continue
                            
                            summaries.append(f"--- Summary of {file_info['name']} ---\n{summary}")

                        elif content_type == 'image':
                            image_parts.append({"name": file_info['name'], "bytes": content})

                        else:
                            st.warning(f"Skipping unsupported file: {file_info['name']} ({content})")

                if summaries or image_parts:
                    # --- REDUCE STEP: Synthesize the final answer ---
                    st.info(f"Synthesizing final answer from {len(summaries)} summaries and {len(image_parts)} images.")
                    with st.spinner("🤖 AI is thinking..."):
                        
                        # Prepare the final prompt context
                        final_context = "Based on the following summaries and images, please answer the user's question.\n\n"
                        if summaries:
                            final_context += "\n\n".join(summaries)
                        
                        # Build the message payload for the final, high-quality model
                        message_content = [{"type": "text", "text": final_context}]
                        message_content.append({"type": "text", "text": f"\n\n--- USER QUESTION ---\n{user_prompt}"})

                        for img in image_parts:
                            base64_image = base64.b64encode(img['bytes']).decode('utf-8')
                            message_content.append({"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}})
                        
                        final_messages = [{"role": "user", "content": message_content}]
                        
                        # Use the powerful model for the final, nuanced answer
                        final_response = get_ai_response(api_key, model="gpt-4o", messages=final_messages)

                        st.markdown("### 🤖 AI Analysis")
                        st.markdown(final_response)

    except HttpError as error:
        st.error(f"A Google Drive API error occurred: {error}")
        st.info("Your Google authentication may have expired. Please try logging out and logging back in.")
